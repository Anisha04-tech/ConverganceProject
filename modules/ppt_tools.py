import os
import platform
import tempfile
from pptx import Presentation
from fpdf import FPDF
from PIL import Image

def convert_ppt(ppt_file, target_ext):
    if target_ext == "pdf":
        try:
            # Try image-based conversion first (requires Windows + PowerPoint)
            return convert_ppt_to_pdf_with_images(ppt_file)
        except Exception as e:
            print(f"[WARNING] Image conversion failed: {e}. Falling back to text-based conversion.")
            return convert_ppt_to_pdf_text(ppt_file)
    else:
        raise ValueError(f"Conversion from PPTX to {target_ext.upper()} is not supported.")

def convert_ppt_to_pdf_text(ppt_file):
    prs = Presentation(ppt_file)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)  # Use built-in font to avoid Unicode font download

    for i, slide in enumerate(prs.slides):
        pdf.add_page()
        text_content = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
            elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text_content.append(para.text.strip())

        if text_content:
            for text in text_content:
                try:
                    pdf.multi_cell(0, 10, text.encode('latin-1', 'replace').decode('latin-1'))
                except Exception as e:
                    pdf.cell(0, 10, "(Text encoding failed)")
        else:
            pdf.cell(0, 10, f"(Slide {i + 1}) No readable text found.")

    output_path = ppt_file.replace(".pptx", "_text.pdf")
    pdf.output(output_path)
    return output_path

def convert_ppt_to_pdf_with_images(ppt_file):
    if platform.system() != "Windows":
        raise EnvironmentError("Image-based slide export only works on Windows with PowerPoint installed.")

    try:
        import comtypes.client
    except ImportError:
        raise ImportError("Missing dependency 'comtypes'. Please run: pip install comtypes")

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(ppt_file, WithWindow=False)

    temp_dir = tempfile.mkdtemp()
    image_folder = os.path.join(temp_dir, "slides")
    os.makedirs(image_folder, exist_ok=True)
    presentation.SaveAs(image_folder, 17)  # 17 = ppSaveAsJPG
    presentation.Close()
    powerpoint.Quit()

    images = sorted([os.path.join(image_folder, f) for f in os.listdir(image_folder) if f.endswith(".JPG")])
    if not images:
        raise RuntimeError("No slide images found. PowerPoint export may have failed.")

    image_list = [Image.open(img).convert("RGB") for img in images]
    pdf_path = ppt_file.replace(".pptx", "_image.pdf")
    image_list[0].save(pdf_path, save_all=True, append_images=image_list[1:])
    return pdf_path
